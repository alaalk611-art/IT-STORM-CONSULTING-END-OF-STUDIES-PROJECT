# app.py — StormCopilot (IT-STORM) — Local HF only (UPDATED)
from __future__ import annotations

import base64
import io
import os
import sys
import time
from pathlib import Path
from typing import List, Optional, Tuple

# =============== Anti-crash: limiter threads & parallelism ===============
os.environ.setdefault("TOKENIZERS_PARALLELISM", "false")
os.environ.setdefault("OMP_NUM_THREADS", "1")
os.environ.setdefault("MKL_NUM_THREADS", "1")
os.environ.setdefault("NUMEXPR_MAX_THREADS", "1")

# --- Import path (project root) ---
sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

import pandas as pd
import plotly.graph_objects as go
import requests
# --- Std / 3rd-party ---
import streamlit as st
from dotenv import load_dotenv
from langchain_core.runnables import RunnablePassthrough
from langchain_community.embeddings import SentenceTransformerEmbeddings
# --- LangChain / RAG ---
from langchain_community.llms import HuggingFacePipeline
from langchain_community.vectorstores import Chroma as LCChroma
from langchain_core.output_parsers import StrOutputParser
from langchain_core.prompts import PromptTemplate

# --- PyTorch (set threads early) ---
try:
    import torch
    torch.set_num_threads(max(1, min(2, (os.cpu_count() or 2))))
except Exception:
    pass

# --- Transformers (local LLM + summarizer) ---
from transformers import AutoModelForCausalLM, AutoTokenizer
from transformers import pipeline as hf_pipeline

try:
    from transformers import AutoModelForSeq2SeqLM as _AM
    from transformers import AutoTokenizer as _AT
    from transformers import pipeline
    HF_AVAILABLE = True
except Exception:
    HF_AVAILABLE = False

# --- Chroma / Embeddings ---
import chromadb
from chromadb.config import Settings
from docx import Document
# --- Office formats ---
from pptx import Presentation
from sentence_transformers import SentenceTransformer

# --- Paths / constants ---
ROOT = Path(__file__).resolve().parents[2]
DATA_RAW = ROOT / "data" / "raw"
DATA_PROCESSED = ROOT / "data" / "processed"
VEC_DIR = ROOT / "vectors"  # garde ton dossier existant si différent
OUT_DIR = ROOT / "out"
COLLECTION_NAME = "itstorm_docs"
EMB_MODEL_NAME = "sentence-transformers/all-MiniLM-L6-v2"

# --- .env ---
env_example = ROOT / ".env.example"
if env_example.exists():
    load_dotenv(dotenv_path=env_example, override=True)
else:
    load_dotenv(override=True)

# --- Page config (unique) ---
st.set_page_config(
    page_title="StormCopilot · IT-STORM",
    page_icon="💼",
    layout="wide",
    initial_sidebar_state="expanded",
)

# --------------------------------------------------------------------------------------
# THEME & BRANDING
# --------------------------------------------------------------------------------------
LOGO_PATH = os.path.join(os.path.dirname(__file__), "assets", "itstorm_logo.png")

THEME_LIGHT = dict(
    bg="#ffffff", text="#0e1626", subbg="#f5f7fb", card="#ffffff", border="#e7edf7", accent="#1a73e8"
)
THEME_DARK = dict(
    bg="#0b1220", text="#e6f0ff", subbg="#111a2b", card="#0f1a2e", border="#1f2b42", accent="#4da3ff"
)

def _logo_b64() -> str:
    try:
        with open(LOGO_PATH, "rb") as f:
            return base64.b64encode(f.read()).decode("utf-8")
    except Exception:
        return ""

def apply_theme(theme: dict):
    css = f"""
    <style>
    :root {{
      --bg: {theme['bg']};
      --text: {theme['text']};
      --subbg: {theme['subbg']};
      --card: {theme['card']};
      --border: {theme['border']};
      --accent: {theme['accent']};
    }}

    * {{ transition: background-color .25s ease, color .25s ease, border-color .25s ease; }}

    html, body, [data-testid="stAppViewContainer"] {{
      background: var(--bg) !important;
      color: var(--text) !important;
    }}
    [data-testid="stSidebar"] {{ background: var(--subbg) !important; }}

    [data-testid="stAppViewContainer"] > .main .block-container {{
      padding-top: 2.6rem !important;
      padding-bottom: 2rem !important;
    }}

    header[data-testid="stHeader"] {{
      background: transparent !important;
      box-shadow: none !important;
    }}

    .brandbar {{
      background: var(--card);
      border: 1px solid var(--border);
      border-radius: 18px;
      padding: 10px 14px;
      margin: .25rem 0 10px;
      box-shadow: 0 8px 24px rgba(0,0,0,.08);
    }}
    .brand-title {{ font-weight: 800; font-size: 22px; line-height: 1.1; }}
    .brand-sub   {{ opacity: .85; font-size: 13px; }}

    .brandbar img {{
      display: block;
      height: 56px;
      width: auto;
      object-fit: contain;
      border-radius: 12px;
      filter: drop-shadow(0 2px 6px rgba(0,0,0,.25));
    }}

    .big-title {{ font-size: 2.0rem; font-weight: 800; color: var(--text); margin-bottom: .25rem; }}
    .sub-title {{ font-size: .95rem; color: var(--text); opacity: .75; }}

    .kpi-card {{
      border: 1px solid var(--border);
      border-radius: 16px;
      padding: 14px 16px;
      background: var(--card);
      box-shadow: 0 1px 2px rgba(10, 30, 80, 0.05);
    }}
    .kpi-label {{ color: var(--text); opacity: .7; font-size: .8rem; margin-bottom: .15rem; }}
    .kpi-value {{ color: var(--text); font-size: 1.35rem; font-weight: 700; }}

    .card {{
      border: 1px solid var(--border);
      border-radius: 16px;
      padding: 18px;
      background: var(--card);
      box-shadow: 0 1px 2px rgba(10, 30, 80, 0.05);
    }}

    .stButton>button, .stDownloadButton>button {{
      border-radius: 10px;
      padding: .55rem .9rem;
      font-weight: 600;
      border: 1px solid var(--border);
      background: var(--accent);
      color: #fff;
    }}
    .stButton>button:hover {{ filter: brightness(1.06); }}

    .chip {{
      display: inline-block;
      padding: 3px 9px;
      background: var(--subbg);
      border: 1px solid var(--border);
      border-radius: 9999px;
      color: var(--text);
      font-size: 12px;
      margin-right: 6px;
    }}
    .source {{
      font-size: 12px;
      color: var(--text);
      opacity: .75;
      background: var(--subbg);
      padding: 4px 8px;
      border-radius: 9999px;
      display: inline-block;
      margin-right: 8px;
      margin-top: 6px;
    }}

    .stTabs [role="tablist"] button {{ border-radius: 10px; }}
    hr {{ border: none; border-top: 1px solid var(--border); margin: .8rem 0 1rem; }}
    </style>
    """
    st.markdown(css, unsafe_allow_html=True)

def render_brand_header(api_ok: Optional[bool]=None, llm_ok: Optional[bool]=None):
    logo = _logo_b64()
    c1, c2, c3 = st.columns([0.10, 0.72, 0.18])
    with c1:
        if logo:
            st.markdown(
                f"""
                <div class="brandbar" style="display:flex;align-items:center;justify-content:center;">
                    <img src="data:image/png;base64,{logo}" width="56" style="border-radius:12px;"/>
                </div>
                """,
                unsafe_allow_html=True
            )
        else:
            st.markdown('<div class="brandbar"><span class="brand-title">IS</span></div>', unsafe_allow_html=True)
    with c2:
        st.markdown(
            """
            <div class="brandbar">
              <div class="brand-title">StormCopilot</div>
              <div class="brand-sub">Hosted by <b>IT-STORM · Innovation & Consulting</b> — Knowledge, RAG & Market Intelligence</div>
            </div>
            """,
            unsafe_allow_html=True,
        )
    with c3:
        st.markdown('<div class="brandbar" style="text-align:right;">', unsafe_allow_html=True)
        if api_ok is not None:
            st.markdown(f'<span class="chip">Market API: {"✅ Up" if api_ok else "❌ Down"}</span>', unsafe_allow_html=True)
        if llm_ok is not None:
            st.markdown(f'<span class="chip">LLM (Local HF): {"✅ Ready" if llm_ok else "⚠️ Not ready"}</span>', unsafe_allow_html=True)
        st.markdown('<span class="chip">Env: Local</span>', unsafe_allow_html=True)
        st.markdown('</div>', unsafe_allow_html=True)

# --------------------------------------------------------------------------------------
# Ensure folders exist
# --------------------------------------------------------------------------------------
for p in [DATA_RAW, DATA_PROCESSED, VEC_DIR, OUT_DIR]:
    p.mkdir(parents=True, exist_ok=True)

# --------------------------------------------------------------------------------------
# CACHED HELPERS (DB, LLM, Retriever, etc.)
# --------------------------------------------------------------------------------------
@st.cache_resource
def get_db():
    client = chromadb.Client(Settings(persist_directory=str(VEC_DIR), is_persistent=True))
    names = [c.name for c in client.list_collections()]
    coll = client.get_collection(COLLECTION_NAME) if COLLECTION_NAME in names else client.create_collection(COLLECTION_NAME)
    return client, coll

@st.cache_resource
def get_embedder():
    return SentenceTransformer(EMB_MODEL_NAME)

@st.cache_resource
def get_summarizer():
    if not HF_AVAILABLE:
        return None
    for name in ["google/pegasus-xsum", "t5-small"]:
        try:
            tok = _AT.from_pretrained(name)
            mdl = _AM.from_pretrained(name)
            return pipeline("summarization", model=mdl, tokenizer=tok)
        except Exception:
            continue
    return None

@st.cache_resource
def get_lc_embeddings():
    return SentenceTransformerEmbeddings(model_name=EMB_MODEL_NAME)

@st.cache_resource
def get_llm(model_name: str = "TinyLlama/TinyLlama-1.1B-Chat-v1.0"):
    """
    Local HF pipeline (CPU par défaut) — aucun appel réseau.
    Ajuste max_new_tokens/temperature ici si besoin.
    Plus robuste: pad/eos définis, low_cpu_mem_usage, dtype auto.
    """
    tok = AutoTokenizer.from_pretrained(model_name)
    if tok.pad_token_id is None and tok.eos_token_id is not None:
        tok.pad_token = tok.eos_token

    mdl = AutoModelForCausalLM.from_pretrained(
        model_name,
        low_cpu_mem_usage=True,
        torch_dtype="auto",
        device_map=None  # CPU
    )

    gen = hf_pipeline(
        task="text-generation",
        model=mdl,
        tokenizer=tok,
        max_new_tokens=256,     # un peu plus court par défaut (stabilité)
        do_sample=False,
        temperature=0.0,
        repetition_penalty=1.1,
        pad_token_id=tok.pad_token_id or tok.eos_token_id,
    )
    return HuggingFacePipeline(pipeline=gen)

@st.cache_resource
def get_retriever():
    lc_emb = get_lc_embeddings()
    vs = LCChroma(
        collection_name=COLLECTION_NAME,
        persist_directory=str(VEC_DIR),
        embedding_function=lc_emb
    )
    return vs.as_retriever(search_type="similarity")

RAG_PROMPT = PromptTemplate.from_template("""
You are a consulting assistant. Answer the user using ONLY the provided context.
If the answer is not in the context, say "Je ne sais pas.".

Answer concisely (<= 10 lines) and include a short bullet list of key points.
Cite sources like: [Source: filename].

Question:
{question}

Context:
{context}

Answer:
""".strip())

def build_rag_chain(k: int = 4):
    retriever = get_retriever()
    llm = get_llm()
    def _format_docs(docs):
        parts = []
        for d in docs:
            src = d.metadata.get("source", "unknown")
            parts.append(f"[Source: {Path(src).name}] {d.page_content}")
        return "\n\n".join(parts)
    chain = (
        {
            "question": RunnablePassthrough(),
            "context": retriever.with_config(run_name="Retriever") | (lambda docs: _format_docs(docs))
        }
        | RAG_PROMPT
        | llm
        | StrOutputParser()
    )
    chain = chain.with_config(configurable={"search_kwargs": {"k": k}})
    return chain

def retrieve(query: str, k: int = 4, similarity_threshold: float = 0.0):
    _, coll = get_db()
    emb = get_embedder()
    q = emb.encode([query], normalize_embeddings=True).tolist()[0]
    res = coll.query(query_embeddings=[q], n_results=k)
    docs = res.get("documents", [[]])[0]
    metas = res.get("metadatas", [[]])[0]
    dists = res.get("distances", [[]])[0] if res.get("distances") else [0.0] * len(docs)
    pairs = []
    for doc, meta, dist in zip(docs, metas, dists):
        if similarity_threshold > 0 and dist > similarity_threshold:
            continue
        pairs.append((doc, meta.get("source", "unknown"), dist))
    return pairs

def save_uploaded_files(files):
    saved = []
    for f in files:
        target = DATA_RAW / f.name
        with open(target, "wb") as out:
            out.write(f.read())
        saved.append(str(target))
    return saved

def extract_and_chunk(in_dir=DATA_RAW, out_path=DATA_PROCESSED / "chunks.jsonl", chunk_size=900, overlap=150):
    import json

    import docx2txt
    from pypdf import PdfReader
    def extract_pdf(p):
      try:
        reader = PdfReader(p)
        return "\n".join(page.extract_text() or "" for page in reader.pages)
      except Exception:
        return ""
    def extract_docx(p):
      try:
        return docx2txt.process(p)
      except Exception:
        return ""
    recs = []
    for p in Path(in_dir).glob("*"):
        text = ""
        if p.suffix.lower() == ".pdf":
            text = extract_pdf(p)
        elif p.suffix.lower() == ".docx":
            text = extract_docx(p)
        elif p.suffix.lower() in {".txt", ".md"}:
            try:
                text = p.read_text(encoding="utf-8", errors="ignore")
            except Exception:
                text = ""
        if not text:
            continue
        words = text.split()
        i = 0
        while i < len(words):
            chunk_words = words[i:i+chunk_size]
            recs.append({"source": str(p), "text": " ".join(chunk_words)})
            i += max(1, (chunk_size - overlap))
    out_path.parent.mkdir(parents=True, exist_ok=True)
    with open(out_path, "w", encoding="utf-8") as f:
        for r in recs:
            f.write(__import__("json").dumps(r, ensure_ascii=False) + "\n")
    return len(recs), out_path

def rebuild_index(in_path=DATA_PROCESSED / "chunks.jsonl"):
    client, _ = get_db()
    try:
        client.delete_collection(COLLECTION_NAME)
    except Exception:
        pass
    coll = client.create_collection(COLLECTION_NAME)
    model = get_embedder()
    texts, metas, ids = [], [], []
    import json
    with open(in_path, "r", encoding="utf-8") as f:
        for k, line in enumerate(f):
            r = json.loads(line)
            texts.append(r["text"])
            metas.append({"source": r["source"]})
            ids.append(f"id_{k}")
    emb = model.encode(texts, normalize_embeddings=True).tolist()
    coll.add(documents=texts, embeddings=emb, metadatas=metas, ids=ids)
    return len(texts)

def kpi_card(label: str, value: str):
    st.markdown(f"""
    <div class="kpi-card">
      <div class="kpi-label">{label}</div>
      <div class="kpi-value">{value}</div>
    </div>
    """, unsafe_allow_html=True)

def render_sources(rows: List[Tuple[str, str, float]]):
    if not rows:
        st.info("No sources retrieved yet.")
        return
    for i, (doc, src, dist) in enumerate(rows, 1):
        with st.expander(f"Context #{i} — {Path(src).name}"):
            st.write(doc[:1200] + ("..." if len(doc) > 1200 else ""))
            st.markdown(f'<span class="source">Source: {src}</span>', unsafe_allow_html=True)

def summarize_text(text: str, max_chars=1200):
    if not text:
        return ""
    text = text[:4000]
    if get_summarizer() is None:
        parts = [p.strip() for p in text.split("\n") if p.strip()]
        return " ".join(parts[:4])[:max_chars]
    summarizer = get_summarizer()
    try:
        out = summarizer(text, max_length=128, min_length=50, do_sample=False)
        return out[0]["summary_text"]
    except Exception:
        return text[:max_chars]

def make_ppt(title="Client Update", bullets=None) -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if bullets:
        s2 = prs.slides.add_slide(prs.slide_layouts[1])
        s2.shapes.title.text = "Key Points"
        tf = s2.shapes.placeholders[1].text_frame
        tf.clear()
        for b in bullets:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0
    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.read()

def make_docx(title="Executive Summary", paragraphs=None) -> bytes:
    doc = Document()
    doc.add_heading(title, 0)
    if paragraphs:
        for p in paragraphs:
            doc.add_paragraph(p)
    bio = io.BytesIO()
    doc.save(bio)
    bio.seek(0)
    return bio.read()

# --------------------------------------------------------------------------------------
# SIDEBAR — THEME TOGGLE (Dark / Light)
# --------------------------------------------------------------------------------------
st.sidebar.header("Settings")
dark_mode = st.sidebar.toggle("Dark mode", value=True, help="Basculer clair/sombre")
apply_theme(THEME_DARK if dark_mode else THEME_LIGHT)

st.sidebar.caption("IT-STORM · Innovation & Consulting")

top_k = st.sidebar.slider("Top-K results", 2, 12, 4, 1)
sim_thresh = st.sidebar.slider("Similarity threshold (distance)", 0.0, 1.5, 0.0, 0.05)
model_name = st.sidebar.text_input("Embedding model", EMB_MODEL_NAME)
if model_name != EMB_MODEL_NAME:
    st.sidebar.warning("This UI uses MiniLM; change in code to switch model safely.")

st.sidebar.markdown("---")
st.sidebar.subheader("Diagnostics")
st.sidebar.markdown("LLM backend: **Local Hugging Face** ✅")
st.sidebar.caption("Theme auto · v2.2 (Local-only, CPU-safe)")

# --------------------------------------------------------------------------------------
# HEADER
# --------------------------------------------------------------------------------------
render_brand_header(api_ok=None, llm_ok=True)

# KPI row
k1, k2, k3, k4 = st.columns(4)
with k1: kpi_card("Indexed Docs", f"{len(list(DATA_RAW.glob('*')))} files")
with k2: kpi_card("Chunks File", "✅" if (DATA_PROCESSED / "chunks.jsonl").exists() else "❌")
with k3: kpi_card("Vector DB", "✅" if any(VEC_DIR.glob('*')) else "❌")
with k4: kpi_card("Output Folder", "✅" if any(OUT_DIR.glob('*')) else "—")

# --------------------------------------------------------------------------------------
# TABS
# --------------------------------------------------------------------------------------
tab1, tab2, tab3, tab4 = st.tabs(["💬 Knowledge Chat", "📂 Upload & Index", "📝 Generate Docs", "🌍 Market Watch"])

# ---- TAB 1: CHAT ----
with tab1:
    st.markdown("### Knowledge Q&A (RAG with LangChain + Local HF)")
    q = st.text_input("Ask a question about your indexed documents")
    btn_search = st.button("🔎 Retrieve & Answer", use_container_width=True)

    k = top_k
    thresh = sim_thresh

    if btn_search and q:
        if len(q) > 2000:
            st.warning("Votre question est très longue — j’ai tronqué à 2000 caractères pour stabilité.")
            q = q[:2000]
        with st.spinner("Retrieving context & generating answer..."):
            try:
                chain = build_rag_chain(k=k)
                answer = chain.invoke(q)
                st.markdown("#### Answer")
                st.write(answer)
                st.markdown("#### Sources (raw chunks)")
                rows = retrieve(q, k=k, similarity_threshold=thresh)
                render_sources(rows)
            except Exception as e:
                st.error(f"RAG chain failed: {e}")
                st.info("Tip: if you switched embedding model, rebuild your index first.")

# ---- TAB 2: UPLOAD & INDEX ----
with tab2:
    st.markdown("### Upload documents & (re)build index")
    files = st.file_uploader("Upload PDF/DOCX/TXT", type=["pdf", "docx", "txt", "md"], accept_multiple_files=True)
    c1, c2 = st.columns(2)
    if c1.button("⬆️ Save uploads", use_container_width=True) and files:
        saved = save_uploaded_files(files)
        st.success(f"Saved {len(saved)} file(s) to `data/raw/`.")
    if c2.button("⚙️ Extract & Chunk", use_container_width=True):
        with st.spinner("Extracting & chunking..."):
            n, outp = extract_and_chunk()
        st.success(f"Created {n} chunks → {outp}")
    if st.button("🧠 Rebuild Vector Index", use_container_width=True):
        with st.spinner("Encoding & indexing... this can take a moment"):
            total = rebuild_index()
        st.success(f"Indexed {total} chunks into ChromaDB.")

    if (DATA_PROCESSED / "chunks.jsonl").exists():
        st.markdown("#### Preview extracted chunks")
        import itertools
        import json
        rows = []
        with open(DATA_PROCESSED / "chunks.jsonl", "r", encoding="utf-8") as f:
            for line in itertools.islice(f, 3):
                rows.append(json.loads(line))
        df = pd.DataFrame(rows)
        st.dataframe(df, use_container_width=True, height=200)

# ---- TAB 3: DOCS GENERATION ----
with tab3:
    st.markdown("### Generate consulting deliverables")
    default_title = "Executive Summary – Client X"
    title = st.text_input("Document Title", value=default_title)
    st.caption("Tip: Paste text below or run a query in the Chat tab and copy the answer here.")
    txt = st.text_area("Content (will be summarized if too long)", height=220)

    c1, c2 = st.columns(2)
    if c1.button("📄 Download DOCX", use_container_width=True):
        bullets = [b.strip() for b in txt.split("\n") if b.strip()]
        paragraphs = bullets if bullets else ["No content provided."]
        data = make_docx(title=title, paragraphs=paragraphs)
        st.download_button("⬇️ Save DOCX", data=data, file_name="report.docx",
                           mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                           use_container_width=True)
    if c2.button("📊 Download PPTX", use_container_width=True):
        bullets = [b.strip() for b in txt.split("\n") if b.strip()]
        data = make_ppt(title=title, bullets=bullets[:8] if bullets else ["No content provided."])
        st.download_button("⬇️ Save PPTX", data=data, file_name="deck.pptx",
                           mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                           use_container_width=True)

# ---- TAB 4: MARKET WATCH ----
with tab4:
    API_BASE = os.getenv("MARKET_API_BASE_URL", "http://127.0.0.1:8001").rstrip("/")

    def api_get(path, **params):
        try:
            r = requests.get(API_BASE + path, params=params, timeout=10)
            r.raise_for_status()
            return r.json()
        except Exception as e:
            st.error(f"API error: {e}")
            return None

    st.header("📈 Market Watch — CAC40")

    symbols = st.multiselect(
        "Choisissez des symboles (Yahoo: suffixe .PA)",
        ["^FCHI", "BNP.PA", "AIR.PA", "MC.PA", "OR.PA", "ORA.PA"],
        ["^FCHI", "BNP.PA", "MC.PA"]
    )

    cols = st.columns(min(len(symbols), 4) or 1)
    for i, sym in enumerate(symbols):
        data = api_get(f"/v1/quote/{sym}")
        with cols[i % len(cols)]:
            if data:
                st.metric(label=sym, value=f"{data.get('price', 0):,.2f} {data.get('currency','')}")
            else:
                st.error(f"Erreur pour {sym}")

    # Historique OHLCV (chandelier)
    st.subheader("Historique (chandelier)")
    chosen = st.selectbox("Historique pour:", options=symbols or ["^FCHI"])
    interval = st.selectbox("Intervalle", options=["1d", "1h", "30m"], index=0)
    rng = st.selectbox("Période", options=["1mo", "3mo", "6mo", "1y"], index=0)

    hist = api_get(f"/v1/ohlcv/{chosen}", interval=interval, range=rng)

    def _pick(df, candidates):
        for name in candidates:
            if name in df.columns:
                return name
        raise KeyError(f"Colonnes manquantes, attendues parmi: {candidates}, trouvé: {list(df.columns)}")

    if hist and hist.get("data"):
        df = pd.DataFrame(hist["data"])
        if df.empty:
            st.info("Pas de données disponibles pour ce symbole/intervalle.")
        else:
            df.columns = [str(c).strip().lower() for c in df.columns]
            ts_col   = _pick(df, ["timestamp", "datetime", "date", "time"])
            open_col = _pick(df, ["open", "o"])
            high_col = _pick(df, ["high", "h"])
            low_col  = _pick(df, ["low", "l"])
            close_col= _pick(df, ["close", "c"])
            vol_col  = next((c for c in ["volume", "vol", "v"] if c in df.columns), None)
            try:
                df[ts_col] = pd.to_datetime(df[ts_col])
            except Exception:
                pass
            fig = go.Figure(data=[go.Candlestick(
                x=df[ts_col], open=df[open_col], high=df[high_col], low=df[low_col], close=df[close_col]
            )])
            st.plotly_chart(fig, use_container_width=True)
            st.dataframe(df[[ts_col, open_col, high_col, low_col, close_col] + ([vol_col] if vol_col else [])].tail(10),
                         use_container_width=True)
    else:
        st.info("Pas de données disponibles pour cce symbole/intervalle.")
