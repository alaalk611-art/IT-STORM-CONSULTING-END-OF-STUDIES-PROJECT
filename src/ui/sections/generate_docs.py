# src/ui/sections/generate_docs.py
from __future__ import annotations
import io
from datetime import datetime
from pathlib import Path
import streamlit as st

# =============== UTILS ===============

def _slugify(name: str) -> str:
    return "".join(c if c.isalnum() or c in ("-", "_") else "_" for c in name).strip("_")

def _now_tag() -> str:
    return datetime.now().strftime("%Y%m%d-%H%M%S")

def build_docx_bytes(title, subtitle, author, date_str, theme_hex, logo_bytes, sections):
    try:
        from docx import Document
        from docx.shared import Inches, Pt
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except Exception as e:
        raise RuntimeError("python-docx n’est pas installé. → pip install python-docx") from e

    doc = Document()

    # Marges
    for s in doc.sections:
        s.top_margin, s.bottom_margin = Inches(0.8), Inches(0.8)
        s.left_margin, s.right_margin = Inches(0.8), Inches(0.8)

    # Logo
    if logo_bytes:
        from tempfile import NamedTemporaryFile
        with NamedTemporaryFile(delete=False, suffix=".png") as tmp:
            tmp.write(logo_bytes)
            tmp.flush()
            p_logo = doc.add_paragraph()
            run_logo = p_logo.add_run()
            run_logo.add_picture(tmp.name, width=Inches(1.2))
            p_logo.alignment = WD_ALIGN_PARAGRAPH.RIGHT

    # Titre principal
    p_title = doc.add_paragraph()
    run = p_title.add_run(title)
    run.bold = True
    run.font.size = Pt(24)

    if subtitle:
        p_sub = doc.add_paragraph(subtitle)
        p_sub.runs[0].font.size = Pt(14)

    meta = doc.add_paragraph(f"{author} — {date_str}")
    meta.runs[0].font.size = Pt(10)

    # Séparation
    doc.add_paragraph("")

    for sec in sections:
        head = sec.get("heading", "").strip()
        body = sec.get("body", "").strip()
        if not head and not body:
            continue
        if head:
            p = doc.add_paragraph(head)
            p.runs[0].bold = True
            p.runs[0].font.size = Pt(14)
        if body:
            for paragraph in body.split("\n\n"):
                doc.add_paragraph(paragraph.strip())

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()

def build_pptx_bytes(title, subtitle, author, date_str, theme_hex, logo_bytes, sections):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except Exception as e:
        raise RuntimeError("python-pptx n’est pas installé. → pip install python-pptx") from e

    def hex_to_rgb(h):
        h = h.replace("#", "")
        if len(h) == 3:
            h = "".join([c * 2 for c in h])
        return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))

    theme_rgb = hex_to_rgb(theme_hex or "#0ea5e9")
    prs = Presentation()
    blank = prs.slide_layouts[6]

    # Slide titre
    s0 = prs.slides.add_slide(blank)
    shape = s0.shapes.add_shape(
        autoshape_type_id=1,
        left=Inches(0), top=Inches(0),
        width=prs.slide_width, height=Inches(1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*theme_rgb)

    # Titre
    tx = s0.shapes.add_textbox(Inches(0.6), Inches(0.2),
                               prs.slide_width - Inches(2), Inches(1.5))
    tf = tx.text_frame
    tf.text = title
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Sous-titre
    tx2 = s0.shapes.add_textbox(Inches(0.6), Inches(1.3),
                                prs.slide_width - Inches(1.2), Inches(1.5))
    tf2 = tx2.text_frame
    tf2.text = f"{subtitle}\n{author} — {date_str}"
    tf2.paragraphs[0].font.size = Pt(18)

    # Logo
    if logo_bytes:
        from tempfile import NamedTemporaryFile
        with NamedTemporaryFile(delete=False, suffix=".png") as tmp:
            tmp.write(logo_bytes)
            tmp.flush()
            s0.shapes.add_picture(tmp.name, prs.slide_width - Inches(1.7),
                                  Inches(0.15), height=Inches(0.7))

    # Slides sections
    for sec in sections:
        head, body = sec.get("heading", ""), sec.get("body", "")
        if not head and not body:
            continue
        slide = prs.slides.add_slide(blank)

        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4),
                                      prs.slide_width - Inches(1.2), Inches(1))
        tf = tb.text_frame
        tf.text = head or "Section"
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(28)

        tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(1.4),
                                       prs.slide_width - Inches(1.2), prs.slide_height - Inches(2))
        tf2 = tb2.text_frame
        for line in body.split("\n"):
            if not line.strip():
                continue
            p = tf2.add_paragraph()
            p.text = f"• {line.strip()}"
            p.font.size = Pt(18)

        foot = slide.shapes.add_textbox(Inches(0.6), prs.slide_height - Inches(0.6),
                                        prs.slide_width - Inches(1.2), Inches(0.4))
        ft = foot.text_frame
        ft.text = date_str
        ft.paragraphs[0].font.size = Pt(10)
        ft.paragraphs[0].alignment = PP_ALIGN.RIGHT

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

# =============== MAIN RENDERER ===============

def render_generate_docs():
    st.header("📝 Generate Docs")

    with st.form("gen_docs_form", clear_on_submit=False):
        colA, colB = st.columns([2, 1])
        with colA:
            title = st.text_input("Titre du document", value="StormCopilot — Rapport & Synthèse")
            subtitle = st.text_input("Sous-titre", value="Projet PFE — IT Storm")
            author = st.text_input("Auteur", value="Ala BEN LAKHAL")
            date_str = st.text_input("Date", value=datetime.now().strftime("%d %B %Y"))
        with colB:
            theme_hex = st.color_picker("Couleur du thème", value="#0EA5E9")
            logo_file = st.file_uploader("Logo (PNG/JPG)", type=["png", "jpg", "jpeg"])
            export_docx = st.checkbox("Exporter DOCX", value=True)
            export_pptx = st.checkbox("Exporter PPTX", value=True)

        st.markdown("#### Sections")
        h1 = st.text_input("Section 1 — Titre", "Résumé exécutif")
        b1 = st.text_area("Section 1 — Contenu", "• Objectif du projet\n• Résultats clés")
        h2 = st.text_input("Section 2 — Titre", "Détails & Méthodologie")
        b2 = st.text_area("Section 2 — Contenu", "• Pipeline RAG\n• LLM & génération locale")
        h3 = st.text_input("Section 3 — Titre", "Conclusions")
        b3 = st.text_area("Section 3 — Contenu", "• Enseignements\n• Étapes suivantes")

        submitted = st.form_submit_button("Générer les documents")

    if submitted:
        if not export_docx and not export_pptx:
            st.warning("Sélectionne au moins un format.")
            return

        logo_bytes = logo_file.read() if logo_file else None
        sections = [
            {"heading": h1, "body": b1},
            {"heading": h2, "body": b2},
            {"heading": h3, "body": b3},
        ]

        base = _slugify(title) or "document"
        tag = _now_tag()

        if export_docx:
            try:
                docx_bytes = build_docx_bytes(title, subtitle, author, date_str, theme_hex, logo_bytes, sections)
                st.download_button(
                    "📄 Télécharger DOCX",
                    data=docx_bytes,
                    file_name=f"{base}_{tag}.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                )
            except Exception as e:
                st.error(f"Erreur DOCX : {e}")

        if export_pptx:
            try:
                pptx_bytes = build_pptx_bytes(title, subtitle, author, date_str, theme_hex, logo_bytes, sections)
                st.download_button(
                    "📊 Télécharger PPTX",
                    data=pptx_bytes,
                    file_name=f"{base}_{tag}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            except Exception as e:
                st.error(f"Erreur PPTX : {e}")
