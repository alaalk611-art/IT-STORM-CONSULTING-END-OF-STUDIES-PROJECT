from __future__ import annotations
# app.py — StormCopilot (IT-STORM) — Local HF only (UPDATED)
# === Chatbot (mutualisé) ===
from src.chatbot import ask_rag, get_chain_cached
from pathlib import Path  # si pas déjà importé
from src.ui.i18n import set_lang_from_query, get_lang, t
import streamlit as st
from src.ui.sections import  auth
from src.ui.sections import speech_chat
from src.ui.sections import tech_watch
from src.ui.sections import home
from src.ui.sections import automation
from src.ui.sections import upload
from src.ui.sections import mlops

import os
os.environ["CHROMA_TELEMETRY_ENABLED"] = "false"



# ⚠️ IMPORTANT : ne le faire qu'une seule fois par session
if "lang_initialized" not in st.session_state:
    set_lang_from_query()          # lit ?lang=fr et initialise st.session_state["lang"]
    st.session_state["lang_initialized"] = True
    # lit ?lang=fr et initialise st.session_state["lang"]
from src.ui.components import floating_chatbot
import streamlit as st
from streamlit_javascript import st_javascript
from src.ui.components.floating_chatbot import render_chatbot
from src.ui.tabs import generate_docs_rag 
from src.utils.doc_tools import make_docx, make_ppt, summarize_text
from src.ui.sections import market

# === UI principale ===
import os, io, base64, time, sys
from src.ui.tabs import generate_docs_rag
from pathlib import Path
from typing import List, Tuple, Optional
import io
from datetime import datetime
from typing import List, Dict
import os
import streamlit as st

# ===================== I18N (EN/FR) =====================
if "lang" not in st.session_state:
    st.session_state["lang"] = "fr"

LANG = {
    
    "fr": {
        "sidebar_settings": "Paramètres",
        "sidebar_dark_mode": "Mode sombre",
        "sidebar_diag": "Diagnostics",
        "sidebar_llm_backend": "Moteur LLM : **Hugging Face local** ✅",
        "sidebar_theme_note": "Thème auto · v2.2 (Local uniquement, CPU-safe)",

        "kpi_indexed_docs": "Documents indexés",
        "kpi_chunks_file": "Fichier de segments",
        "kpi_vector_db": "Base vectorielle",
        "kpi_output_folder": "Dossier d’export",

        "tab_chat": "💬 Chat Connaissance",
        "tab_upload": "📂 Import & Index",
        "tab_generate": "📝 Générer des documents",
        "tab_market": "🌍 Marchés",
        "tab_qa": "🔎 QA anglais (RAG EN)",

        "chat_title": "Q&R Connaissance (RAG avec LangChain + HF local)",
        "chat_input": "Pose une question sur tes documents indexés",
        "chat_button": "🔎 Récupérer & Répondre",
        "chat_clear": "🧹 Effacer",   # ✅ ajouté ici
        "chat_answer": "Réponse",
        "chat_sources": "Sources (extraits bruts)",
        "chat_question_too_long": "Ta question est très longue — je l’ai tronquée à 2000 caractères pour la stabilité.",
        "chat_tip_rebuild": "Astuce : si tu as changé de modèle d’embedding, reconstruis l’index avant.",

        "upload_title": "Importer des documents & (re)construire l’index",
        "upload_uploader": "Importer PDF/DOCX/TXT",
        "upload_save": "⬆️ Enregistrer les fichiers",
        "upload_saved": "✅ {n} fichier(s) enregistré(s) dans `data/raw/`.",
        "upload_extract": "⚙️ Extraire & Segmenter",
        "upload_extracted": "Créé {n} segments → {outp}",
        "upload_rebuild": "🧠 Reconstruire l’index vectoriel",
        "upload_rebuilt": "Indexé {n} segments dans ChromaDB.",
        "upload_preview": "Aperçu des segments extraits",

        "theme_env": "Env : Local",
        "brand_market_api": "API Marché : ",
        "brand_llm": "LLM (HF local) : ",
        "brand_up": "✅ OK",
        "brand_down": "❌ Hors ligne",
        "brand_ready": "✅ Prêt",
        "brand_not_ready": "⚠️ Non prêt",

        "lang_picker": "Langue",
        
        "lang_fr": "Français",
    },
}

def t(key: str) -> str:
    lang = st.session_state.get("lang", "fr")
    return LANG.get(lang, LANG["fr"]).get(key, key)
# ============================================================

from src.rag.chain import build_rag_chain
from src.ui.tabs import generate_docs_rag  # <— ce fichier


# Optionnel : valeurs par défaut si .env non chargé
os.environ.setdefault("LLM_MODEL", "google/flan-t5-base")
os.environ.setdefault("LLM_DEVICE", "-1")
os.environ.setdefault("EMBEDDING_MODEL", "sentence-transformers/all-MiniLM-L6-v2")
os.environ.setdefault("RAG_TOPK", "4")
os.environ.setdefault("RAG_MAX_CONTEXT_CHARS", "1300")
def _get_source_names_for_query(query: str, k: int = 4) -> List[str]:
    try:
        rows = retrieve(query, k=k, similarity_threshold=0.0)  # (doc, src, dist)
        names = []
        for _doc, src, _dist in rows:
            try:
                names.append(Path(src).name)
            except Exception:
                names.append(str(src))
        # unique en conservant l'ordre
        seen = set()
        uniq = []
        for n in names:
            if n not in seen:
                seen.add(n)
                uniq.append(n)
        return uniq
    except Exception:
        return []

@st.cache_resource(show_spinner=False)
def load_chain():
    return build_rag_chain()

# =============== Anti-crash: limiter threads & parallelism ===============
os.environ.setdefault("TOKENIZERS_PARALLELISM", "false")
os.environ.setdefault("OMP_NUM_THREADS", "1")
os.environ.setdefault("MKL_NUM_THREADS", "1")
os.environ.setdefault("NUMEXPR_MAX_THREADS", "1")

# --- Import path (project root) ---
sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

# --- Std / 3rd-party ---
import streamlit as st
import pandas as pd
import requests
import plotly.graph_objects as go
from dotenv import load_dotenv
API_BASE = os.getenv("MARKET_API_BASE_URL", "http://127.0.0.1:8001")
TIMEOUT = float(os.getenv("FINANCE_TIMEOUT", "10"))
# --- LangChain / RAG ---
from langchain_community.llms import HuggingFacePipeline
from langchain_community.vectorstores import Chroma as LCChroma
from langchain_community.embeddings import SentenceTransformerEmbeddings
from langchain_core.prompts import PromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.runnables import RunnablePassthrough

# --- PyTorch (set threads early) ---
try:
    import torch
    torch.set_num_threads(max(1, min(2, (os.cpu_count() or 2))))
except Exception:
    pass

# --- Transformers (local LLM + summarizer) ---
from transformers import AutoModelForCausalLM, AutoTokenizer, pipeline as hf_pipeline
try:
    from transformers import pipeline, AutoTokenizer as _AT, AutoModelForSeq2SeqLM as _AM
    HF_AVAILABLE = True
except Exception:
    HF_AVAILABLE = False

# --- Chroma / Embeddings ---
import chromadb
from chromadb.config import Settings
from sentence_transformers import SentenceTransformer

# --- Office formats ---
from pptx import Presentation
from docx import Document

# --- Paths / constants ---
ROOT = Path(__file__).resolve().parents[2]
DATA_RAW = ROOT / "data" / "raw"
DATA_PROCESSED = ROOT / "data" / "processed"
VEC_DIR = ROOT / "vectors"  # garde ton dossier existant si différent
OUT_DIR = ROOT / "out"
COLLECTION_NAME = "itstorm_docs"
EMB_MODEL_NAME = "sentence-transformers/all-MiniLM-L6-v2"

# --- .env ---
env_example = ROOT / ".env.example"
if env_example.exists():
    load_dotenv(dotenv_path=env_example, override=True)
else:
    load_dotenv(override=True)

# --- Page config (unique) ---
st.set_page_config(
    page_title="StormCopilot · IT-STORM",
    page_icon="💼",
    layout="wide",
    initial_sidebar_state="expanded",
)
# ============================================================
# UI Layout optimization — centrer et réduire marges
# ============================================================
st.markdown("""
<style>
/* Supprime les marges bleues et recentre tout le contenu principal */
header[data-testid="stHeader"] {
    height: 0 !important;
    visibility: hidden !important;
}

/* APRÈS — flux normal, rien ne “mange” le haut de page */
.block-container {
    /* revenir au flux normal */
    display: block;
    min-height: initial;
    padding-top: 2rem !important;   /* petite marge en haut */
    padding-bottom: 1.5rem !important;
    margin: 0 auto !important;
}

section.main {
    display: block;
}

/* si tu veux cacher l’header, compense son retrait par un padding-top global */
header[data-testid="stHeader"] {
    height: 0 !important;
    visibility: hidden !important;
}
.stAppViewContainer .main .block-container {
    padding-top: 3.5rem !important; /* compense l’header caché */
}


/* Supprime les marges / espaces inutiles en haut et bas */
.stAppViewContainer, .main {
    padding: 0 !important;
    margin: 0 !important;
}

/* Cache les diviseurs, si visibles */
hr, [role="separator"] {
    display: none !important;
}
</style>
""", unsafe_allow_html=True)

# src/api/main.py
import os
from typing import List, Any, Dict

from fastapi import FastAPI, HTTPException, Query, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, HTMLResponse
from fastapi.staticfiles import StaticFiles

from dotenv import load_dotenv

# === Tes imports existants (marchés) ==========================
from src.api.models import Quote, OHLCV, SearchResult
from src.api.providers.yahoo import YahooProvider
from src.api.utils import is_euronext_open, paris_now
# ===============================================================

load_dotenv()

app = FastAPI(
    title="IT-Storm API",
    version="1.1.0",
    description="Market API + Chat bubble endpoint (RAG) pour it-storm.fr",
)

# --- CORS ---
def _parse_origins(env_val: str) -> List[str]:
    if not env_val:
        return []
    return [o.strip() for o in env_val.split(",") if o.strip()]

# Par défaut on autorise it-storm.fr + localhost dev ; surcharge possible par CORS_ORIGINS
_default_origins = [
    "https://it-storm.fr",
    "https://www.it-storm.fr",
    "http://localhost:3000",
    "http://127.0.0.1:3000",
    "http://localhost:8000",
    "http://127.0.0.1:8000",
]
allow_origins = _parse_origins(os.getenv("CORS_ORIGINS")) or _default_origins

app.add_middleware(
    CORSMiddleware,
    allow_origins=allow_origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- Static: servir la bulle (embed.js) ---
# Dossier vu dans ton arbo : client_chatbot/public/embed.js
STATIC_CLIENT_DIR = os.path.join(os.getcwd(), "client_chatbot", "public")
if os.path.isdir(STATIC_CLIENT_DIR):
    app.mount("/client-chat", StaticFiles(directory=STATIC_CLIENT_DIR), name="client-chat")
else:
    # On ne crashe pas l'appli si le dossier n'existe pas ; simple warning en logs
    print(f"[WARN] Static dir not found for chat bubble: {STATIC_CLIENT_DIR}")

# =======================
#  ROUTES MARCHÉ (existantes)
# =======================
PROVIDER = YahooProvider()

@app.get("/v1/health")
async def health():
    return {"status": "ok", "time": paris_now().isoformat()}

@app.get("/v1/marketstatus")
async def market_status():
    return {"market": "Euronext Paris", "is_open": is_euronext_open(), "time": paris_now().isoformat()}

@app.get("/v1/search", response_model=List[SearchResult])
async def search(q: str = Query(..., min_length=1)):
    return await PROVIDER.search(q)

@app.get("/v1/quote/{symbol}", response_model=Quote)
async def get_quote(symbol: str):
    try:
        return await PROVIDER.quote(symbol)
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))

@app.get("/v1/ohlcv/{symbol}", response_model=OHLCV)
async def get_ohlcv(symbol: str, interval: str = "1m", range: str = "1d"):
    try:
        return await PROVIDER.ohlcv(symbol, interval=interval, range_=range)
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))

# =======================
#  CHAT / RAG (nouveaux)
# =======================

# Chargement paresseux de la chaîne QA (même logique que tools/qa_cli)
_QA_CHAIN: Any = None
_QA_LOADED: bool = False

def _try_build_chain() -> Any:
    """
    Essaie plusieurs méthodes classiques pour récupérer la même chaîne
    utilisée par `tools/qa_cli.py`.
    Adapte si besoin selon tes vrais noms de fonctions.
    """
    # 1) src.rag.chain: get_chain / build_chain
    try:
        from src.rag.chain import get_chain  # type: ignore
        return get_chain()
    except Exception:
        pass
    try:
        from src.rag.chain import build_chain  # type: ignore
        return build_chain()
    except Exception:
        pass

    # 2) src.rag (ex: get_chain au niveau du package)
    try:
        from src import rag  # type: ignore
        if hasattr(rag, "get_chain"):
            return rag.get_chain()
        if hasattr(rag, "build_chain"):
            return rag.build_chain()
    except Exception:
        pass

    # 3) tools.qa_cli: chaîne exposée ?
    try:
        # Si ton qa_cli expose une fonction utilitaire (à adapter si tu en as une)
        from tools.qa_cli import get_chain as cli_get_chain  # type: ignore
        return cli_get_chain()
    except Exception:
        pass

    # 4) Rien trouvé
    raise RuntimeError(
        "Impossible de charger la chaîne QA. "
        "Expose une fonction `get_chain()` ou `build_chain()` dans `src.rag.chain` "
        "ou adapte `_try_build_chain()` à ton projet."
    )

def _ensure_chain() -> Any:
    global _QA_CHAIN, _QA_LOADED
    if _QA_LOADED and _QA_CHAIN is not None:
        return _QA_CHAIN
    _QA_CHAIN = _try_build_chain()
    _QA_LOADED = True
    return _QA_CHAIN

def _invoke_chain(chain: Any, query: str) -> str:
    """
    Normalise l'appel pour différents types de chaînes (LCEL LangChain, fonctions, etc.).
    """
    # Cas LangChain LCEL: chain.invoke({"query": ...})
    try:
        out = chain.invoke({"query": query})
        if isinstance(out, dict):
            # Champs possibles: 'result', 'answer', 'output_text', etc.
            for k in ("result", "answer", "output_text", "text", "content"):
                if k in out and out[k]:
                    return str(out[k])
            # Sinon, tout le dict
            return str(out)
        return str(out)
    except Exception:
        pass

    # Cas fonction simple: chain(query)
    try:
        return str(chain(query))
    except Exception:
        pass

    # Cas méthode .run
    try:
        return str(chain.run(query))
    except Exception:
        pass

    # Dernier recours
    return "Je n’ai pas pu générer de réponse (chaîne non compatible)."

@app.post("/chat")
async def chat(req: Request):
    """
    Endpoint utilisé par la bulle JS (embed.js) ou tout frontend.
    Reçoit { "message": "..." } et renvoie { "answer", "sources", "suggestions" }.
    """
    try:
        payload: Dict[str, Any] = await req.json()
    except Exception:
        raise HTTPException(status_code=400, detail="Payload JSON invalide")

    q = str(payload.get("message") or "").strip()
    if not q:
        raise HTTPException(status_code=422, detail="Champ 'message' requis")

    try:
        # k = top_k UI si dispo, sinon défaut 4
        k = 4
        try:
            # si un slider Streamlit a posé top_k dans session_state (facultatif)
            k = int(st.session_state.get("top_k_from_ui", 4))
        except Exception:
            pass

        result = ask_rag(
            question=q,
            k=k,
            get_sources_fn=_get_source_names_for_query
        )
        return JSONResponse(result)
    except Exception as e:
        return JSONResponse(
            {"answer": f"[Erreur backend] {e}", "sources": [], "suggestions": []},
            status_code=500
        )

# Petite page de test locale (facultatif)
@app.get("/")
async def home():
    return HTMLResponse("""
<!doctype html>
<html lang="fr">
  <head>
    <meta charset="utf-8">
    <title>IT-Storm — Test bulle</title>
  </head>
  <body>
    <h1>Test bulle IT-Storm</h1>
    <p>Si la bulle apparaît en bas à droite, clique et pose une question.</p>

    <script defer
      src="/client-chat/embed.js"
      data-api="/chat"
      data-brand="#2563eb"
      data-title="IT-Storm Chatbot"></script>
  </body>
</html>
""" )

# --------------------------------------------------------------------------------------
# THEME & BRANDING
# --------------------------------------------------------------------------------------
LOGO_PATH = os.path.join(os.path.dirname(__file__), "assets", "itstorm_logo.png")

THEME_LIGHT = dict(
    bg="#ffffff", text="#0e1626", subbg="#f5f7fb", card="#ffffff", border="#e7edf7", accent="#1a73e8"
)
THEME_DARK = dict(
    bg="#0b1220", text="#e6f0ff", subbg="#111a2b", card="#0f1a2e", border="#1f2b42", accent="#4da3ff"
)

def _logo_b64() -> str:
    try:
        with open(LOGO_PATH, "rb") as f:
            return base64.b64encode(f.read()).decode("utf-8")
    except Exception:
        return ""

def apply_theme(theme: dict):
    css = f"""
    <style>
    :root {{
      --bg: {theme['bg']};
      --text: {theme['text']};
      --subbg: {theme['subbg']};
      --card: {theme['card']};
      --border: {theme['border']};
      --accent: {theme['accent']};
    }}

    /* =========================
       BASE
       ========================= */
    * {{
      transition:
        background-color .25s ease,
        color .25s ease,
        border-color .25s ease;
    }}

    html, body, [data-testid="stAppViewContainer"] {{
      background: var(--bg) !important;
      color: var(--text) !important;
    }}

    [data-testid="stSidebar"] {{
      background: var(--subbg) !important;
    }}

    [data-testid="stAppViewContainer"] > .main .block-container {{
      padding-top: 2.6rem !important;
      padding-bottom: 2rem !important;
    }}

    header[data-testid="stHeader"] {{
      background: transparent !important;
      box-shadow: none !important;
    }}

    /* =========================
       BRAND / CARTES
       ========================= */
    .brandbar {{
      background: var(--card);
      border: 1px solid var(--border);
      border-radius: 18px;
      padding: 10px 14px;
      margin: .25rem 0 10px;
      box-shadow: 0 8px 24px rgba(0,0,0,.08);
    }}
    .brand-title {{
      font-weight: 800;
      font-size: 22px;
      line-height: 1.1;
    }}
    .brand-sub {{
      opacity: .85;
      font-size: 13px;
    }}

    .brandbar img {{
      display: block;
      height: 56px;
      width: auto;
      object-fit: contain;
      border-radius: 12px;
      filter: drop-shadow(0 2px 6px rgba(0,0,0,.25));
    }}

    .big-title {{
      font-size: 2.0rem;
      font-weight: 800;
      color: var(--text);
      margin-bottom: .25rem;
    }}
    .sub-title {{
      font-size: .95rem;
      color: var(--text);
      opacity: .75;
    }}

    .kpi-card {{
      border: 1px solid var(--border);
      border-radius: 16px;
      padding: 14px 16px;
      background: var(--card);
      box-shadow: 0 1px 2px rgba(10, 30, 80, 0.05);
    }}
    .kpi-label {{
      color: var(--text);
      opacity: .7;
      font-size: .8rem;
      margin-bottom: .15rem;
    }}
    .kpi-value {{
      color: var(--text);
      font-size: 1.35rem;
      font-weight: 700;
    }}

    .card {{
      border: 1px solid var(--border);
      border-radius: 16px;
      padding: 18px;
      background: var(--card);
      box-shadow: 0 1px 2px rgba(10, 30, 80, 0.05);
    }}

    /* =========================
       BOUTONS
       ========================= */
    .stButton>button,
    .stDownloadButton>button {{
      border-radius: 10px;
      padding: .55rem .9rem;
      font-weight: 600;
      border: 1px solid var(--border);
      background: var(--accent);
      color: #fff;
    }}
    .stButton>button:hover,
    .stDownloadButton>button:hover {{
      filter: brightness(1.06);
    }}

    /* =========================
       CHIPS / SOURCES
       ========================= */
    .chip {{
      display: inline-block;
      padding: 3px 9px;
      background: var(--subbg);
      border: 1px solid var(--border);
      border-radius: 9999px;
      color: var(--text);
      font-size: 12px;
      margin-right: 6px;
    }}
    .source {{
      font-size: 12px;
      color: var(--text);
      opacity: .75;
      background: var(--subbg);
      padding: 4px 8px;
      border-radius: 9999px;
      display: inline-block;
      margin-right: 8px;
      margin-top: 6px;
    }}

    /* =========================
       NAVIGATION PRINCIPALE
       GRID 4 colonnes (2 lignes) — centrée — pills plus larges
       + onglet actif illuminé + barre lumineuse
       ========================= */

    /* Wrapper : centre la nav */
    div[data-testid="stRadio"] {{
      display: flex !important;
      justify-content: center !important;
      padding-left: 8.5rem !important;
      padding-right: 3.5rem !important;
    }}

    
    /* Radiogroup = GRID */
    div[data-testid="stRadio"] > div[role="radiogroup"] {{
      display: grid !important;
      grid-template-columns: repeat(4, minmax(0, 1fr)) !important;

      gap: 0.85rem !important;
      align-items: center !important;
      justify-items: stretch !important;

      max-width: 1050px;
      width: 100%;
      margin: 0 auto !important;

      padding: 0.45rem 0.25rem 0.6rem 0.25rem !important;
    }}

    /* Cache le rond radio */
    div[data-testid="stRadio"] > div[role="radiogroup"] > label > div:first-child {{
      display: none !important;
    }}

    /* Force pleine largeur (wrappers internes Streamlit) */
    div[data-testid="stRadio"] > div[role="radiogroup"] > label,
    div[data-testid="stRadio"] > div[role="radiogroup"] > label > div {{
      width: 100% !important;
      box-sizing: border-box !important;
    }}

    /* Pills */
    div[data-testid="stRadio"] > div[role="radiogroup"] > label {{
      position: relative !important;
      display: flex !important;
      justify-content: center !important;
      align-items: center !important;

      min-height: 44px;

      border-radius: 999px !important;
      border: 1px solid rgba(226,232,240,0.95) !important;
      background: rgba(248,250,252,0.96) !important;

      padding: 0.5rem 1.25rem !important;
      font-size: 0.95rem !important;
      font-weight: 650 !important;
      color: #1e293b !important;

      box-shadow: 0 2px 6px rgba(15,23,42,0.06) !important;
      cursor: pointer !important;
      user-select: none !important;

      transition: transform .18s ease, box-shadow .18s ease,
                  border-color .18s ease, background .18s ease !important;
    }}

    /* Hover */
    div[data-testid="stRadio"] > div[role="radiogroup"] > label:hover {{
      background: rgba(238,242,255,0.97) !important;
      border-color: rgba(148,163,184,0.9) !important;
      transform: translateY(-1px) !important;
      box-shadow: 0 12px 26px rgba(2,6,23,0.10) !important;
    }}

    /* Barre lumineuse (off) */
    div[data-testid="stRadio"] > div[role="radiogroup"] > label::after {{
      content: "";
      position: absolute;
      left: 18%;
      right: 18%;
      bottom: -8px;
      height: 4px;
      border-radius: 999px;
      opacity: 0;
      transform: translateY(-2px);
      background: transparent;
      transition: opacity .18s ease, transform .18s ease,
                  box-shadow .18s ease, background .18s ease !important;
      pointer-events: none;
    }}

    /* Onglet actif (illuminé) */
    div[data-testid="stRadio"] > div[role="radiogroup"] > label[data-checked="true"] {{
      background: linear-gradient(135deg, #ffffff, #e0f2fe) !important;
      border-color: rgba(56,189,248,1) !important;
      color: #075985 !important;

      transform: translateY(-2px) !important;
      box-shadow:
        0 16px 36px rgba(56,189,248,0.30),
        0 0 0 4px rgba(56,189,248,0.18) !important;
    }}

    /* Barre lumineuse active */
    div[data-testid="stRadio"] > div[role="radiogroup"] > label[data-checked="true"]::after {{
      opacity: 1;
      transform: translateY(0);
      background: linear-gradient(90deg,
        rgba(56,189,248,0.0),
        rgba(56,189,248,1),
        rgba(99,102,241,0.95),
        rgba(56,189,248,0.0)
      );
      box-shadow:
        0 12px 28px rgba(56,189,248,0.45),
        0 0 22px rgba(56,189,248,0.35);
    }}

    /* Responsive */
    @media (max-width: 900px) {{
      div[data-testid="stRadio"] > div[role="radiogroup"] {{
        grid-template-columns: repeat(2, 1fr) !important;
        max-width: 680px;
      }}
    }}

    @media (max-width: 520px) {{
      div[data-testid="stRadio"] > div[role="radiogroup"] {{
        grid-template-columns: 1fr !important;
        max-width: 100%;
      }}
    }}

    /* Ligne de séparation douce */
    hr {{
      border: none;
      border-top: 1px solid var(--border);
      margin: .8rem 0 1rem;
    }}
    </style>
    """
    st.markdown(css, unsafe_allow_html=True)


def render_brand_header(api_ok: Optional[bool]=None, llm_ok: Optional[bool]=None):
    logo = _logo_b64()
    c1, c2 = st.columns([0.14, 0.86])

    with c1:
        if logo:
            st.markdown(
                f"""
                <div class="brandbar" style="height: 78px; display:flex; align-items:center; justify-content:center;">
                    <img src="data:image/png;base64,{logo}" width="56" style="border-radius:12px;"/>
                </div>
                """,
                unsafe_allow_html=True
            )
        else:
            st.markdown(
                '<div class="brandbar" style="height: 78px; display:flex; align-items:center; justify-content:center;">'
                '<span class="brand-title">IS</span></div>',
                unsafe_allow_html=True
            )

    with c2:
        st.markdown(
            """
            <div class="brandbar" style="height: 78px; display:flex; flex-direction:column; justify-content:center;">
              <div class="brand-title">StormCopilot</div>
              <div class="brand-sub">Hosted by <b>IT-STORM · Innovation & Consulting</b> — Knowledge, RAG & Market Intelligence</div>
            </div>
            """,
            unsafe_allow_html=True,
        )
 

# --------------------------------------------------------------------------------------
# Ensure folders exist
# --------------------------------------------------------------------------------------
for p in [DATA_RAW, DATA_PROCESSED, VEC_DIR, OUT_DIR]:
    p.mkdir(parents=True, exist_ok=True)

# --------------------------------------------------------------------------------------
# CACHED HELPERS (DB, LLM, Retriever, etc.)
# --------------------------------------------------------------------------------------
@st.cache_resource
def get_db():
    client = chromadb.Client(Settings(persist_directory=str(VEC_DIR), is_persistent=True))
    names = [c.name for c in client.list_collections()]
    coll = client.get_collection(COLLECTION_NAME) if COLLECTION_NAME in names else client.create_collection(COLLECTION_NAME)
    return client, coll

@st.cache_resource
def get_embedder():
    return SentenceTransformer(EMB_MODEL_NAME)

@st.cache_resource
def get_summarizer():
    if not HF_AVAILABLE:
        return None
    for name in ["google/pegasus-xsum", "t5-small"]:
        try:
            tok = _AT.from_pretrained(name)
            mdl = _AM.from_pretrained(name)
            return pipeline("summarization", model=mdl, tokenizer=tok)
        except Exception:
            continue
    return None

@st.cache_resource
def get_lc_embeddings():
    return SentenceTransformerEmbeddings(model_name=EMB_MODEL_NAME)

@st.cache_resource
def get_llm(model_name: str = "TinyLlama/TinyLlama-1.1B-Chat-v1.0"):
    """
    Local HF pipeline (CPU par défaut) — aucun appel réseau.
    Ajuste max_new_tokens/temperature ici si besoin.
    Plus robuste: pad/eos définis, low_cpu_mem_usage, dtype auto.
    """
    tok = AutoTokenizer.from_pretrained(model_name)
    if tok.pad_token_id is None and tok.eos_token_id is not None:
        tok.pad_token = tok.eos_token

    mdl = AutoModelForCausalLM.from_pretrained(
        model_name,
        low_cpu_mem_usage=True,
        torch_dtype="auto",
        device_map=None  # CPU
    )

    gen = hf_pipeline(
        task="text-generation",
        model=mdl,
        tokenizer=tok,
        max_new_tokens=256,     # un peu plus court par défaut (stabilité)
        do_sample=False,
        temperature=0.0,
        repetition_penalty=1.1,
        pad_token_id=tok.pad_token_id or tok.eos_token_id,
    )
    return HuggingFacePipeline(pipeline=gen)

@st.cache_resource
def get_retriever():
    lc_emb = get_lc_embeddings()
    vs = LCChroma(
        collection_name=COLLECTION_NAME,
        persist_directory=str(VEC_DIR),
        embedding_function=lc_emb
    )
    return vs.as_retriever(search_type="similarity")

RAG_PROMPT = PromptTemplate.from_template("""
You are a consulting assistant. Answer the user using ONLY the provided context.
If the answer is not in the context, say "Je ne sais pas.".

Answer concisely (<= 10 lines) and include a short bullet list of key points.
Cite sources like: [Source: filename].

Question:
{question}

Context:
{context}

Answer:
""".strip())

def build_rag_chain(k: int = 4):
    retriever = get_retriever()
    llm = get_llm()
    def _format_docs(docs):
        parts = []
        for d in docs:
            src = d.metadata.get("source", "unknown")
            parts.append(f"[Source: {Path(src).name}] {d.page_content}")
        return "\n\n".join(parts)
    chain = (
        {
            "question": RunnablePassthrough(),
            "context": retriever.with_config(run_name="Retriever") | (lambda docs: _format_docs(docs))
        }
        | RAG_PROMPT
        | llm
        | StrOutputParser()
    )
    chain = chain.with_config(configurable={"search_kwargs": {"k": k}})
    return chain

def retrieve(query: str, k: int = 4, similarity_threshold: float = 0.0):
    _, coll = get_db()
    emb = get_embedder()
    q = emb.encode([query], normalize_embeddings=True).tolist()[0]
    res = coll.query(query_embeddings=[q], n_results=k)
    docs = res.get("documents", [[]])[0]
    metas = res.get("metadatas", [[]])[0]
    dists = res.get("distances", [[]])[0] if res.get("distances") else [0.0] * len(docs)
    pairs = []
    for doc, meta, dist in zip(docs, metas, dists):
        if similarity_threshold > 0 and dist > similarity_threshold:
            continue
        pairs.append((doc, meta.get("source", "unknown"), dist))
    return pairs

def save_uploaded_files(files):
    saved = []
    for f in files:
        target = DATA_RAW / f.name
        with open(target, "wb") as out:
            out.write(f.read())
        saved.append(str(target))
    return saved

def extract_and_chunk(
    in_dir=DATA_RAW,
    out_path=DATA_PROCESSED / "chunks.jsonl",
    chunk_size=900,
    overlap=150,
):
    """
    Extrait le texte de plusieurs formats (PDF, DOCX, TXT/MD, PPTX, CSV, XLS/XLSX, HTML/HTM, JSON),
    le découpe en segments, puis écrit un JSONL pour l’indexation RAG.
    Retourne (nb_chunks, out_path, stats_dict).
    """
    import json
    from pathlib import Path

    # ---- Helpers par type ----
    def extract_pdf(p: Path) -> str:
        try:
            from pypdf import PdfReader
            reader = PdfReader(p)
            return "\n".join((page.extract_text() or "") for page in reader.pages)
        except Exception:
            return ""

    def extract_docx(p: Path) -> str:
        try:
            import docx2txt
            return docx2txt.process(p)
        except Exception:
            return ""

    def extract_txt_like(p: Path) -> str:
        try:
            return p.read_text(encoding="utf-8", errors="ignore")
        except Exception:
            return ""

    def extract_pptx(p: Path) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(p)
            parts = []
            for slide in prs.slides:
                for sh in slide.shapes:
                    if hasattr(sh, "text") and sh.text:
                        parts.append(sh.text)
            return "\n".join(parts)
        except Exception:
            return ""

    def extract_csv(p: Path) -> str:
        try:
            import pandas as pd
            df = pd.read_csv(p)
            return df.to_string(index=False)
        except Exception:
            return ""

    def extract_xlsx(p: Path) -> str:
        try:
            import pandas as pd
            xls = pd.ExcelFile(p)
            parts = []
            for sheet in xls.sheet_names:
                df = xls.parse(sheet)
                parts.append(f"# Sheet: {sheet}\n{df.to_string(index=False)}")
            return "\n\n".join(parts)
        except Exception:
            return ""

    def extract_html(p: Path) -> str:
        try:
            from bs4 import BeautifulSoup
            html = p.read_text(encoding="utf-8", errors="ignore")
            return BeautifulSoup(html, "html.parser").get_text(separator="\n")
        except Exception:
            # Fallback minimal au cas où
            import re
            html = p.read_text(encoding="utf-8", errors="ignore")
            text = re.sub(r"<(script|style)[^>]*>.*?</\\1>", "", html, flags=re.S|re.I)
            text = re.sub(r"<[^>]+>", " ", text)
            return re.sub(r"\s+", " ", text).strip()

    def extract_json(p: Path) -> str:
        try:
            raw = json.loads(p.read_text(encoding="utf-8", errors="ignore"))
        except Exception:
            return ""
        parts = []
        def walk(x):
            if isinstance(x, dict):
                for k, v in x.items():
                    parts.append(str(k))
                    walk(v)
            elif isinstance(x, list):
                for it in x:
                    walk(it)
            else:
                s = str(x)
                if s and s != "None":
                    parts.append(s)
        walk(raw)
        return "\n".join(parts)

    # ---- Boucle fichiers ----
    recs = []
    stats = {"seen": 0, "processed": 0, "ignored": 0, "by_ext": {}}
    for p in Path(in_dir).glob("*"):
        if not p.is_file():
            continue
        stats["seen"] += 1
        ext = p.suffix.lower()
        stats["by_ext"].setdefault(ext, 0)

        text = ""
        if ext == ".pdf":
            text = extract_pdf(p)
        elif ext == ".docx":
            text = extract_docx(p)
        elif ext in {".txt", ".md"}:
            text = extract_txt_like(p)
        elif ext == ".pptx":
            text = extract_pptx(p)
        elif ext == ".csv":
            text = extract_csv(p)
        elif ext in {".xlsx", ".xls"}:
            text = extract_xlsx(p)
        elif ext in {".html", ".htm"}:
            text = extract_html(p)
        elif ext == ".json":
            text = extract_json(p)
        else:
            stats["ignored"] += 1
            continue

        if not text:
            stats["ignored"] += 1
            continue

        # Chunking
        words = text.split()
        step = max(1, (chunk_size - overlap))
        i = 0
        while i < len(words):
            chunk_words = words[i : i + chunk_size]
            recs.append({"source": str(p), "text": " ".join(chunk_words)})
            i += step

        stats["processed"] += 1
        stats["by_ext"][ext] += 1

    # Écriture JSONL
    out_path.parent.mkdir(parents=True, exist_ok=True)
    with open(out_path, "w", encoding="utf-8") as f:
        for r in recs:
            f.write(json.dumps(r, ensure_ascii=False) + "\n")

    return len(recs), out_path, stats

def rebuild_index(in_path=DATA_PROCESSED / "chunks.jsonl"):
    client, _ = get_db()
    try:
        client.delete_collection(COLLECTION_NAME)
    except Exception:
        pass
    coll = client.create_collection(COLLECTION_NAME)
    model = get_embedder()
    texts, metas, ids = [], [], []
    import json
    with open(in_path, "r", encoding="utf-8") as f:
        for k, line in enumerate(f):
            r = json.loads(line)
            texts.append(r["text"])
            metas.append({"source": r["source"]})
            ids.append(f"id_{k}")
    emb = model.encode(texts, normalize_embeddings=True).tolist()
    coll.add(documents=texts, embeddings=emb, metadatas=metas, ids=ids)
    return len(texts)

def kpi_card(label: str, value: str):
    st.markdown(f"""
    <div class="kpi-card">
      <div class="kpi-label">{label}</div>
      <div class="kpi-value">{value}</div>
    </div>
    """, unsafe_allow_html=True)

def render_sources(rows: List[Tuple[str, str, float]]):
    if not rows:
        st.info("No sources retrieved yet.")
        return
    for i, (doc, src, dist) in enumerate(rows, 1):
        with st.expander(f"Context #{i} — {Path(src).name}"):
            st.write(doc[:1200] + ("..." if len(doc) > 1200 else ""))
            st.markdown(f'<span class="source">Source: {src}</span>', unsafe_allow_html=True)

def summarize_text(text: str, max_chars=1200):
    if not text:
        return ""
    text = text[:4000]
    if get_summarizer() is None:
        parts = [p.strip() for p in text.split("\n") if p.strip()]
        return " ".join(parts[:4])[:max_chars]
    summarizer = get_summarizer()
    try:
        out = summarizer(text, max_length=128, min_length=50, do_sample=False)
        return out[0]["summary_text"]
    except Exception:
        return text[:max_chars]

def make_ppt(title="Client Update", bullets=None) -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if bullets:
        s2 = prs.slides.add_slide(prs.slide_layouts[1])
        s2.shapes.title.text = "Key Points"
        tf = s2.shapes.placeholders[1].text_frame
        tf.clear()
        for b in bullets:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0
    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.read()

def make_docx(title="Executive Summary", paragraphs=None) -> bytes:
    doc = Document()
    doc.add_heading(title, 0)
    if paragraphs:
        for p in paragraphs:
            doc.add_paragraph(p)
    bio = io.BytesIO()
    doc.save(bio)
    bio.seek(0)
    return bio.read()

# --------------------------------------------------------------------------------------
# SIDEBAR — LANGUE, THEME, CARTES INFO (NON FONCTIONNELLES)
# --------------------------------------------------------------------------------------
import base64
from pathlib import Path

# ========= CSS sidebar (flags + cartes + rotateur) =========
st.sidebar.markdown("""
<style>
section[data-testid="stSidebar"] > div {
    padding-top: 1rem;
}

/* Cartes */
.sidebar-card {
    border-radius: 16px;
    padding: 0.9rem 0.85rem 0.8rem 0.85rem;
    margin-bottom: 0.8rem;
    background: rgba(255,255,255,0.55);
    box-shadow: 0 4px 14px rgba(15,23,42,0.05);
    border: 1px solid rgba(148,163,184,0.35);
}
[data-theme="dark"] .sidebar-card {
    background: rgba(15,23,42,0.95);
    border-color: rgba(148,163,184,0.35);
}

/* Titres */
.sidebar-title {
    font-size: 0.9rem;
    font-weight: 700;
    margin-bottom: 0.4rem;
    display: flex;
    align-items: center;
    gap: 0.35rem;
}
.sidebar-title span.icon {
    font-size: 1.0rem;
}

/* Ligne fine */
.sidebar-sep {
    height: 1px;
    background: linear-gradient(90deg,
                                rgba(148,163,184,0.2),
                                rgba(79,70,229,0.45),
                                rgba(148,163,184,0.2));
    margin: 0.15rem 0 0.55rem 0;
    border-radius: 999px;
}


/* Texte d’aide */
.sidebar-help {
    font-size: 0.73rem;
    opacity: 0.8;
    margin-top: 0.25rem;
}

/* Listes */
.sidebar-list {
    font-size: 0.8rem;
    padding-left: 1.1rem;
    margin-bottom: 0.1rem;
}
.sidebar-list li { margin-bottom: 0.15rem; }

/* ------- Rotateur IT-STORM (5 messages, fade in/out) ------- */
.itstorm-rotator {
    position: relative;
    overflow: hidden;
    min-height: 52px;      /* hauteur min du bloc */
}
.itstorm-rotator-inner {
    position: relative;
}

/* Tous les items sont superposés, un seul visible à la fois */
.itstorm-rotator .itstorm-rotator-item {
    position: absolute;
    top: 0;
    left: 0;
    right: 0;
    opacity: 0;
    display: flex;
    align-items: center;
    font-size: 0.8rem;
    line-height: 1.1rem;
    padding-right: 4px;
}

/* Animation de base : visible 0–16%, invisible ensuite */
@keyframes itstormFade {
    0%, 16%   { opacity: 1; }
    20%, 100% { opacity: 0; }
}

/* 5 messages, chacun avec un décalage de 10s sur 50s */
.itstorm-rotator .itstorm-rotator-item:nth-child(1) {
    animation: itstormFade 50s infinite;
}
.itstorm-rotator .itstorm-rotator-item:nth-child(2) {
    animation: itstormFade 50s infinite;
    animation-delay: 10s;
}
.itstorm-rotator .itstorm-rotator-item:nth-child(3) {
    animation: itstormFade 50s infinite;
    animation-delay: 20s;
}
.itstorm-rotator .itstorm-rotator-item:nth-child(4) {
    animation: itstormFade 50s infinite;
    animation-delay: 30s;
}
.itstorm-rotator .itstorm-rotator-item:nth-child(5) {
    animation: itstormFade 50s infinite;
    animation-delay: 40s;
}
</style>
""", unsafe_allow_html=True)


# ========= Helpers langue / flags =========
flag_fr_path = Path("src/ui/assets/flag_fr.png")

def _b64(p: Path) -> str:
    with open(p, "rb") as f:
        return base64.b64encode(f.read()).decode("utf-8")

url_lang = st.query_params.get("lang")
if "lang" not in st.session_state:
    st.session_state["lang"] = url_lang if url_lang in {"en", "fr"} else "en"
else:
    if url_lang in {"en", "fr"} and url_lang != st.session_state["lang"]:
        st.session_state["lang"] = url_lang

lang = st.session_state["lang"]

# ========= CARTE 0 : IT-STORM EN BREF (MICRO-INDICATEURS ROTATIFS) =========
st.sidebar.markdown(
    """
<style>
/* =========================
   ROTATEUR IT-STORM (micro-indicateurs)
   ========================= */

/* 1 ligne, pas de wrap, hauteur stable */
.itstorm-rotator{
  position: relative;
  overflow: hidden;
  height: 34px;
  margin-top: 6px;
}

.itstorm-rotator-inner{
  position: relative;
  height: 34px;
}

/* Un seul visible */
.itstorm-rotator-item{
  position: absolute;
  inset: 0;
  height: 34px;

  display: flex;
  align-items: center;

  opacity: 0;
  white-space: nowrap;
  overflow: hidden;
  text-overflow: ellipsis;

  font-size: 0.78rem;
  font-weight: 600;
  line-height: 1;
  will-change: opacity, transform;
  pointer-events: none;
}

/* 5s avec GAP anti-chevauchement */
@keyframes itstormFade5s{
  0%   { opacity: 0; transform: translateY(6px); }
  14%  { opacity: 1; transform: translateY(0); }
  68%  { opacity: 1; transform: translateY(0); }
  86%  { opacity: 0; transform: translateY(-4px); }
  100% { opacity: 0; transform: translateY(-4px); }
}

/* 6 indicateurs => 30s total */
.itstorm-rotator-item:nth-child(1){ animation: itstormFade5s 30s infinite; animation-delay: 0s;  }
.itstorm-rotator-item:nth-child(2){ animation: itstormFade5s 30s infinite; animation-delay: 5s;  }
.itstorm-rotator-item:nth-child(3){ animation: itstormFade5s 30s infinite; animation-delay: 10s; }
.itstorm-rotator-item:nth-child(4){ animation: itstormFade5s 30s infinite; animation-delay: 15s; }
.itstorm-rotator-item:nth-child(5){ animation: itstormFade5s 30s infinite; animation-delay: 20s; }
.itstorm-rotator-item:nth-child(6){ animation: itstormFade5s 30s infinite; animation-delay: 25s; }


/* =========================
   ROTATEUR TITRES (Navigation / Modules / Paramètres / À propos)
   ========================= */

.sc-title-rotator{
  position: relative;
  overflow: hidden;
  height: 24px;
  line-height: 24px;
}

.sc-title-item{
  position: absolute;
  inset: 0;
  opacity: 0;
  pointer-events: none;

  display: flex;
  align-items: center;

  white-space: nowrap;
  overflow: hidden;
  text-overflow: ellipsis;

  font-weight: 800;
  font-size: 0.92rem;
  will-change: opacity, transform;
}

/* 5s avec GAP anti-chevauchement */
@keyframes scTitleFade5s{
  0%   { opacity: 0; transform: translateY(6px); }
  14%  { opacity: 1; transform: translateY(0); }
  68%  { opacity: 1; transform: translateY(0); }
  86%  { opacity: 0; transform: translateY(-4px); }
  100% { opacity: 0; transform: translateY(-4px); }
}

/* 4 titres => 20s total */
.sc-title-item:nth-child(1){ animation: scTitleFade5s 20s infinite; animation-delay: 0s;  }
.sc-title-item:nth-child(2){ animation: scTitleFade5s 20s infinite; animation-delay: 5s;  }
.sc-title-item:nth-child(3){ animation: scTitleFade5s 20s infinite; animation-delay: 10s; }
.sc-title-item:nth-child(4){ animation: scTitleFade5s 20s infinite; animation-delay: 15s; }


/* =========================
   Accessibilité : réduire animations
   ========================= */
@media (prefers-reduced-motion: reduce){
  .itstorm-rotator-item, .sc-title-item{
    animation: none !important;
    opacity: 1 !important;
    transform: none !important;
  }
}
</style>

<!-- ========= CARTE 0 : IT-STORM EN BREF ========= -->
<div class="sidebar-card">
  <div class="sidebar-title">
    <span class="icon">⚡</span>
    <span>IT-STORM en quelques mots</span>
  </div>
  <div class="sidebar-sep"></div>

  <p style="font-size:0.82rem; line-height:1.25rem; margin: 0 0 6px 0;">
    <strong>StormCopilot</strong> est un copilote IA pour consultants IT-STORM.
  </p>

  <div class="itstorm-rotator">
    <div class="itstorm-rotator-inner">
      <div class="itstorm-rotator-item">📂 Upload & Index : base documentaire</div>
      <div class="itstorm-rotator-item">🧠 RAG : réponses sourcées</div>
      <div class="itstorm-rotator-item">🌍 Market Watch : signaux & tendances</div>
      <div class="itstorm-rotator-item">🔎 Veille techno : Cloud · Data · IA</div>
      <div class="itstorm-rotator-item">⚙️ Automation Studio : workflows n8n</div>
      <div class="itstorm-rotator-item">📊 MLOps : dérive & performance</div>
    </div>
  </div>
</div>



    """,
    unsafe_allow_html=True,
)

# ========= LANGUE =========

# ========= DARK / LIGHT MODE =========
dark_mode_key = "ui_dark_mode_toggle_unique"
if "dark_mode" not in st.session_state:
    st.session_state["dark_mode"] = False

st.sidebar.markdown('<div class="sidebar-card">', unsafe_allow_html=True)
st.sidebar.markdown(
    f'<div class="sidebar-title"><span class="icon">💡</span>'
    f'<span>{"Appearance" if lang=="en" else "Apparence"}</span></div>'
    '<div class="sidebar-sep"></div>',
    unsafe_allow_html=True
)

dark_mode = st.sidebar.toggle(
    "Dark mode" if lang == "en" else "Mode sombre",
    value=st.session_state["dark_mode"],
    help="Toggle light / dark theme" if lang == "en" else "Basculer entre thème clair et sombre",
    key=dark_mode_key,
)
st.session_state["dark_mode"] = dark_mode
apply_theme(THEME_DARK if dark_mode else THEME_LIGHT)

st.sidebar.caption("IT-STORM · Innovation & Consulting")
st.sidebar.markdown("</div>", unsafe_allow_html=True)

# =========================================================
# SIDEBAR — CARTES (FR-only) — Mise en valeur des modules
# =========================================================

# ========= CARTE 1 : BIEN UTILISER STORMCOPILOT =========
st.sidebar.markdown('<div class="sidebar-card">', unsafe_allow_html=True)
st.sidebar.markdown(
    """
    <div class="sidebar-title"><span class="icon">💬</span>
      <span>Bien utiliser StormCopilot</span>
    </div>
    <div class="sidebar-sep"></div>
    """,
    unsafe_allow_html=True
)

st.sidebar.markdown(
    """
<ul class="sidebar-list">
  <li><strong>Commence simple</strong> : une question claire ou un objectif court.</li>
  <li><strong>Choisis ton module</strong> : Documents, Marchés, Voix, Veille, Automations ou MLOps.</li>
  <li><strong>Reste concret</strong> : noms de missions, contexte IT-STORM, attentes client, contraintes.</li>
  <li><strong>Améliore la précision</strong> : si la réponse est trop générale, ajoute 1 détail (périmètre, délai, livrable).</li>
</ul>
    """,
    unsafe_allow_html=True,
)
st.sidebar.markdown("</div>", unsafe_allow_html=True)


# ========= CARTE 2 : MODULES PRINCIPAUX (mise en avant) =========
st.sidebar.markdown('<div class="sidebar-card">', unsafe_allow_html=True)
st.sidebar.markdown(
    """
    <div class="sidebar-title"><span class="icon">🧩</span>
      <span>Modules disponibles</span>
    </div>
    <div class="sidebar-sep"></div>
    """,
    unsafe_allow_html=True
)

st.sidebar.markdown(
    """
<ul class="sidebar-list">
  <li><strong>📂 Upload & Index</strong> : ajouter des documents internes, les préparer et enrichir la base de connaissance.</li>
  <li><strong>📝 Génération de documents</strong> : produire des livrables structurés (résumés, notes, contenus prêts à partager).</li>
  <li><strong>🌍 Market Watch</strong> : suivre les marchés, indicateurs, signaux et tendances utiles à la prise de décision.</li>
  <li><strong>🎤 Assistant vocal</strong> : poser une question à l’oral et obtenir une réponse rapide, orientée mission.</li>
  <li><strong>🔎 Veille technologique</strong> : suivre Cloud, Data, IA, DevOps et tendances du freelancing/consulting.</li>
  <li><strong>⚙️ Automation Studio</strong> : lancer et orchestrer des workflows n8n pour automatiser les routines (radars, refresh, rapports).</li>
  <li><strong>🧠 MLOps Market</strong> : surveiller modèles, qualité, dérive et performance, avec logique champion/challenger.</li>
</ul>
    """,
    unsafe_allow_html=True,
)
st.sidebar.markdown("</div>", unsafe_allow_html=True)


# ========= CARTE 3 : À PROPOS (jury-friendly) =========
st.sidebar.markdown('<div class="sidebar-card">', unsafe_allow_html=True)
st.sidebar.markdown(
    """
    <div class="sidebar-title"><span class="icon">🧾</span>
      <span>À propos de cette version</span>
    </div>
    <div class="sidebar-sep"></div>
    """,
    unsafe_allow_html=True
)

st.sidebar.markdown(
    """
    <div style="font-weight:800; margin-bottom:0.35rem;">
      StormCopilot • Studio IT-STORM
    </div>

<ul class="sidebar-list">
  <li><strong>Version</strong> : <em>Aperçu 2025</em></li>
  <li><strong>Mode</strong> : local, priorité au hors-ligne, orienté fiabilité</li>
  <li><strong>Positionnement</strong> : copilote IA pour consultants IT-STORM, missions et portage salarial</li>
  <li><strong>But</strong> : gagner du temps sur la recherche, la veille, l’automatisation et la préparation de livrables</li>
</ul>
    """,
    unsafe_allow_html=True,
)
st.sidebar.markdown("</div>", unsafe_allow_html=True)

# --------------------------------------------------------------------------------------
# HEADER
# --------------------------------------------------------------------------------------
render_brand_header(api_ok=None, llm_ok=True)

# KPI row
#k1, k2, k3, k4 = st.columns(4)
#with k1:
#    kpi_card(t("kpi_indexed_docs"), f"{len(list(DATA_RAW.glob('*')))} files")
#with k2:
#    kpi_card(t("kpi_chunks_file"), "✅" if (DATA_PROCESSED / "chunks.jsonl").exists() else "❌")
#with k3:
#    kpi_card(t("kpi_vector_db"), "✅" if any(VEC_DIR.glob('*')) else "❌")
#with k4:
#    kpi_card(t("kpi_output_folder"), "✅" if any(OUT_DIR.glob('*')) else "—")

# --------------------------------------------------------------------------------------
# AUTH GATE GLOBAL — page de login avant l'accès aux onglets
# --------------------------------------------------------------------------------------
# --------------------------------------------------------------------------------------
# AUTH GATE GLOBAL — une seule fois par session
# --------------------------------------------------------------------------------------
#if "auth_ok" not in st.session_state:
#    st.session_state["auth_ok"] = False

#if not st.session_state["auth_ok"]:
#    auth_ok = auth.render_auth_gate()
#    if auth_ok:
#        st.session_state["auth_ok"] = True
#        st.rerun()
#    else:
#        st.stop()

# =====================================================================
# NAVIGATION / TABS (radio + URL + liens ouvrables via clic droit)
# =====================================================================

# ---------------------------
# Définition des routes
# ---------------------------
# =====================================================================
# NAVIGATION / TABS (UNIQUE) — liens (clic gauche = même onglet, clic droit = nouvel onglet)
# =====================================================================
# =====================================================================
# NAVIGATION / TABS (UNIQUE) — liens (clic gauche = même onglet, clic droit = nouvel onglet)
# =====================================================================

ROUTES = {
    "home": "🏠 Accueil",
    "upload": "📂 Upload & Index",
    "docs": "📝 Generate Docs",
    "voice": "🎤 Voice Copilot",
    "tech": "🔎 Veille Techno",
    "market": "🌍 Market Watch",
    "mlops": "🧠 MLOps Market",
    "automation": "⚙️ Automation Studio",
    
}

def render_tabs_as_links(routes: dict, current_page: str):
    st.markdown(
        """
        <style>
          .sc-tabs { display:grid; grid-template-columns:repeat(4,minmax(0,1fr));
                    gap:0.85rem; max-width:1050px; margin:0.6rem auto 0.9rem auto; }
          .sc-tab  { text-decoration:none; text-align:center; display:flex; justify-content:center; align-items:center;
                    min-height:44px; border-radius:999px; border:1px solid rgba(226,232,240,0.95);
                    background:rgba(248,250,252,0.96); padding:0.5rem 1.25rem; font-size:0.95rem; font-weight:650;
                    color:#1e293b; box-shadow:0 2px 6px rgba(15,23,42,0.06);
                    transition:transform .18s ease, box-shadow .18s ease, border-color .18s ease, background .18s ease; }
          .sc-tab:hover { background:rgba(238,242,255,0.97); border-color:rgba(148,163,184,0.9);
                         transform:translateY(-1px); box-shadow:0 12px 26px rgba(2,6,23,0.10); }
          .sc-tab.active { background:linear-gradient(135deg,#ffffff,#e0f2fe); border-color:rgba(56,189,248,1);
                          color:#075985; transform:translateY(-2px);
                          box-shadow:0 16px 36px rgba(56,189,248,0.30), 0 0 0 4px rgba(56,189,248,0.18); }
          @media (max-width:900px){ .sc-tabs{ grid-template-columns:repeat(2,1fr); max-width:680px; } }
          @media (max-width:520px){ .sc-tabs{ grid-template-columns:1fr; max-width:100%; } }
        </style>
        """,
        unsafe_allow_html=True,
    )

    items = []
    for key, label in routes.items():
        cls = "sc-tab active" if key == current_page else "sc-tab"
        # IMPORTANT: target="_self" => clic gauche même onglet
        items.append(f'<a class="{cls}" href="?page={key}" target="_self" rel="noopener">{label}</a>')

    st.markdown(f'<div class="sc-tabs">{"".join(items)}</div>', unsafe_allow_html=True)

# ---------------------------
# Lire l’URL actuelle
# ---------------------------
params = st.query_params
current_page = params.get("page", "home")
if isinstance(current_page, list):
    current_page = current_page[0]
if current_page not in ROUTES:
    current_page = "home"

# Afficher navigation unique
render_tabs_as_links(ROUTES, current_page=current_page)

# ---------------------------
# ROUTING (comme tu as déjà)
# ---------------------------
from src.ui.sections import home

if current_page == "home":
    home.render_home_tab()
    render_chatbot()

elif current_page == "upload":
    upload.render_upload_tab()

elif current_page == "docs":
    generate_docs_rag.render()

elif current_page == "market":
    market.render()

elif current_page == "voice":
    speech_chat.render_stt_only()

elif current_page == "tech":
    tech_watch.render()

elif current_page == "automation":
    automation.render_automation_tab()

elif current_page == "mlops":
    mlops.render_mlops_tab()
