import os

from docx import Document
from pptx import Presentation

os.makedirs("out", exist_ok=True)

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Exemple Rapport Projet"
prs.save("out/exemple.pptx")

doc = Document()
doc.add_heading("Compte-rendu Exemple", 0)
doc.add_paragraph("Ceci est un test de génération automatique.")
doc.save("out/exemple.docx")

print("📁 Fichiers générés dans le dossier 'out'")
