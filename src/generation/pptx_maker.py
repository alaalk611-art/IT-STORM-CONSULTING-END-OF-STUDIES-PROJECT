import io

from pptx import Presentation


def make_ppt(title="Client Update", bullets=None) -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if bullets:
        s2 = prs.slides.add_slide(prs.slide_layouts[1])
        s2.shapes.title.text = "Key Points"
        tf = s2.shapes.placeholders[1].text_frame
        tf.clear()
        for b in bullets:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0
    import io

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.read()
